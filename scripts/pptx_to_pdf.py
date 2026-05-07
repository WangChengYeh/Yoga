#!/usr/bin/env python3
"""Convert YogaFlow3D-Proposal.pptx to PDF using python-pptx + reportlab."""
import sys
from pptx import Presentation
from reportlab.lib.pagesizes import landscape, A4
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib.units import cm
from reportlab.pdfgen import canvas
from reportlab.lib.colors import HexColor

BG = HexColor("#111827")
BLUE = HexColor("#7ECFFF")
GOLD = HexColor("#FFDD88")
WHITE = HexColor("#FFFFFF")
GRAY = HexColor("#9CA3AF")
PAGE_W, PAGE_H = landscape(A4)
MARGIN = 1.8 * cm

def draw_slide_bg(c, slide_num, total):
    c.setFillColor(BG)
    c.rect(0, 0, PAGE_W, PAGE_H, fill=1, stroke=0)
    # slide number
    c.setFillColor(GRAY)
    c.setFont("Helvetica", 8)
    c.drawRightString(PAGE_W - MARGIN, MARGIN * 0.5, f"{slide_num} / {total}")

def extract_slide_content(slide):
    """Return (header, title, subtitle, bullets, table_rows) from a slide."""
    texts = []
    table_data = None
    for shape in slide.shapes:
        if shape.has_table:
            tbl = shape.table
            rows = []
            for r in tbl.rows:
                rows.append([c.text_frame.text.strip() for c in r.cells])
            table_data = rows
        elif shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                txt = para.text.strip()
                if txt:
                    texts.append((para.level, txt))

    # texts[0] is usually "YogaFlow 3D" header (skip if duplicate)
    header = texts[0][1] if texts else ""
    remaining = texts[1:] if (texts and texts[0][1] == "YogaFlow 3D") else texts
    title = remaining[0][1] if remaining else header
    rest = remaining[1:] if remaining else []
    subtitle = rest[0][1] if rest else ""
    bullets = rest[1:] if rest else []
    return header, title, subtitle, bullets, table_data

def render_slide(c, slide, slide_num, total):
    draw_slide_bg(c, slide_num, total)
    _, title, subtitle, bullets, table_data = extract_slide_content(slide)

    y = PAGE_H - MARGIN

    # Title bar accent line
    c.setStrokeColor(BLUE)
    c.setLineWidth(2)
    c.line(MARGIN, y - 2, PAGE_W - MARGIN, y - 2)

    # Title
    c.setFillColor(BLUE)
    c.setFont("Helvetica-Bold", 22)
    y -= 1.2 * cm
    # wrap long titles
    max_w = PAGE_W - 2 * MARGIN
    c.drawString(MARGIN, y, title[:90])
    y -= 0.5 * cm

    # Subtitle / quote
    if subtitle:
        c.setFillColor(GOLD)
        c.setFont("Helvetica-Oblique", 13)
        y -= 0.3 * cm
        c.drawString(MARGIN, y, subtitle[:100])
        y -= 0.6 * cm

    y -= 0.2 * cm

    # Table
    if table_data:
        col_count = max(len(r) for r in table_data)
        col_w = (PAGE_W - 2 * MARGIN) / col_count
        row_h = 0.7 * cm
        for ri, row in enumerate(table_data):
            for ci, cell in enumerate(row):
                x = MARGIN + ci * col_w
                if ri == 0:
                    c.setFillColor(BLUE)
                    c.setFont("Helvetica-Bold", 11)
                    c.rect(x, y - row_h, col_w - 4, row_h, fill=1, stroke=0)
                    c.setFillColor(BG)
                else:
                    c.setFillColor(HexColor("#1F2937") if ri % 2 == 0 else HexColor("#111827"))
                    c.rect(x, y - row_h, col_w - 4, row_h, fill=1, stroke=0)
                    c.setFillColor(WHITE)
                c.setFont("Helvetica-Bold" if ri == 0 else "Helvetica", 10)
                c.drawString(x + 6, y - row_h + 8, str(cell)[:45])
            y -= row_h
        y -= 0.4 * cm

    # Bullets
    for level, text in bullets:
        if y < MARGIN + 1 * cm:
            break
        indent = MARGIN + level * 0.6 * cm
        bullet_x = indent
        text_x = indent + 0.4 * cm
        if level == 0:
            c.setFillColor(BLUE)
            c.setFont("Helvetica", 10)
            c.drawString(bullet_x, y, "•")
            c.setFillColor(WHITE)
            c.setFont("Helvetica", 11)
        else:
            c.setFillColor(GOLD)
            c.setFont("Helvetica", 9)
            c.drawString(bullet_x, y, "›")
            c.setFillColor(GRAY)
            c.setFont("Helvetica", 10)
        c.drawString(text_x, y, text[:95])
        y -= 0.65 * cm

def main():
    pptx_path = "docs/YogaFlow3D-Proposal.pptx"
    pdf_path = "docs/YogaFlow3D-Proposal.pdf"

    prs = Presentation(pptx_path)
    slides = list(prs.slides)
    total = len(slides)

    c = canvas.Canvas(pdf_path, pagesize=(PAGE_W, PAGE_H))
    for i, slide in enumerate(slides):
        render_slide(c, slide, i + 1, total)
        c.showPage()

    c.save()
    import os
    size_kb = os.path.getsize(pdf_path) // 1024
    print(f"Done: {pdf_path} ({total} pages, {size_kb} KB)")

if __name__ == "__main__":
    main()
